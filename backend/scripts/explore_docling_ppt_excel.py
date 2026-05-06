"""
Goal: Compare extraction quality between (Docling) and specialised libraries for PPT and Excel files. 
  - Docling (document converter for PDF/DOCX/PPTX/XLSX)
  - python-pptx (PowerPoint specialised library)
  - openpyxl (Excel specialised library)
Output: Extracted text files saved in _local_uploads/ for manual comparison
Usage:
  1. Place test files in: backend/_local_uploads/
  2. Run: python explore_docling_ppt_excel.py
  3. Compare extracted_*.txt files to evaluate quality
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

# ------------------------------
# SECTION 1: PowerPoint Extraction with python-pptx
# ------------------------------

def explore_pptx_extraction(file_path: str):
    """Extract all slide text using python-pptx."""
    print("\n" + "="*80)
    print("SECTION 1: PowerPoint (PPTX) with python-pptx")
    print("="*80)
    
    try:
        from pptx import Presentation
        print("[OK] python-pptx available")
    except ImportError:
        print("[ERROR] python-pptx not installed")
        return
    
    try:
        prs = Presentation(file_path)
        print(f"\nFile: {Path(file_path).name}")
        print(f"Total Slides: {len(prs.slides)}")
        
        all_text = []
        
        # Iterate through all slides to extract text and table content
        for slide_num, slide in enumerate(prs.slides, start=1):
            print(f"\n--- Slide {slide_num} ---")
            slide_text = []
            
            # Process all shapes on slide (text boxes, tables, images, etc.)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_text.append(text)
                    print(f"  {text[:100]}{'...' if len(text) > 100 else ''}")
                
                # Handle tables separately since they need cell-by-cell extraction
                if shape.has_table:
                    print(f"  [Table: {len(shape.table.rows)} rows]")
                    for row in shape.table.rows:
                        # Format table rows as pipe-separated values for readability
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        if row_text.strip():
                            slide_text.append(row_text)
                            print(f"    {row_text}")
            
            if slide_text:
                # Format with markdown headers for better readability in output file
                all_text.append(f"## Slide {slide_num}\n" + "\n".join(slide_text))
        
        full_text = "\n\n".join(all_text)
        print(f"\nTotal characters: {len(full_text)}")
        
        return full_text
        
    except FileNotFoundError:
        print(f"[ERROR] File not found: {file_path}")
        print("Note: Place your .pptx file in: backend/_local_uploads/")
        return None
    except Exception as e:
        print(f"[ERROR] Error processing file: {e}")
        return None


# ------------------------------
# SECTION 2: Excel Extraction with openpyxl
# ------------------------------

def explore_excel_extraction(file_path: str):
    """Extract data from Excel using openpyxl."""
    print("\n" + "="*80)
    print("SECTION 2: Excel (XLSX)")
    print("="*80)
    
    print("\n--- Method 1: openpyxl ---")
    try:
        import openpyxl
        print("[OK] openpyxl available")
        
        wb = openpyxl.load_workbook(file_path, data_only=True)
        print(f"\nFile: {Path(file_path).name}")
        print(f"Sheets: {wb.sheetnames}")
        
        all_text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            print(f"\n--- Sheet: {sheet_name} ---")
            
            sheet_text = [f"## Sheet: {sheet_name}"]
            
            row_count = 0
            # Extract all rows with at least one non-empty cell
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    # Format as pipe-separated values for consistency with PPTX tables
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    sheet_text.append(row_text)
                    row_count += 1
                    

            print(f"  Total rows: {row_count}")
            all_text.append("\n".join(sheet_text))
        
        full_text_openpyxl = "\n\n".join(all_text)
        print(f"\nTotal characters (openpyxl): {len(full_text_openpyxl)}")
        return full_text_openpyxl
        
    except ImportError:
        print("[ERROR] openpyxl not installed. Install with: pip install openpyxl")
        return None
    except FileNotFoundError:
        print(f"[ERROR] File not found: {file_path}")
        return None
    except Exception as e:
        print(f"[ERROR] {e}")
        return None





# ------------------------------
# SECTION 4: Docling PowerPoint Test
# Test Docling's unified converter on PPTX files
# ------------------------------

def test_docling_pptx(file_path: str):
    """Test Docling's PPTX processing - outputs structured markdown."""
    print("\n" + "="*80)
    print("DOCLING PPTX TEST")
    print("="*80)
    
    try:
        from docling.document_converter import DocumentConverter
        print("[OK] Docling available")
        
        if not Path(file_path).exists():
            print(f"[SKIPPED] File not found: {file_path}")
            return None
        
        print(f"\nProcessing: {Path(file_path).name}")
        
        converter = DocumentConverter()
        result = converter.convert(file_path)
        markdown_text = result.document.export_to_markdown()
        
        # Add filename as main header and remove empty image markers
        markdown_text = markdown_text.replace("<!-- image -->", "")
        markdown_text = f"# {Path(file_path).name}\n\n" + markdown_text
        
        print(f"[SUCCESS] Extracted {len(markdown_text)} characters")
        return markdown_text
        
    except ImportError:
        print("[ERROR] Docling not installed")
        return None
    except Exception as e:
        print(f"[ERROR] {e}")
        return None

# ------------------------------
# SECTION 5: Docling Excel Test
# Test Docling's unified converter on XLSX files
# ------------------------------

def test_docling_excel(file_path: str):
    """Test Docling's Excel processing - outputs structured markdown."""
    print("\n" + "="*80)
    print("DOCLING EXCEL TEST")
    print("="*80)
    
    try:
        from docling.document_converter import DocumentConverter
        print("[OK] Docling available")
        
        if not Path(file_path).exists():
            print(f"[SKIPPED] File not found: {file_path}")
            return None
        
        print(f"\nProcessing: {Path(file_path).name}")
        
        converter = DocumentConverter()
        result = converter.convert(file_path)
        markdown_text = result.document.export_to_markdown()
        
        # Add filename as main header and inject Excel sheet names inline
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheet_names = wb.sheetnames
            
            # Remove Docling's <!-- image --> markers (they mark sheet boundaries)
            markdown_text = markdown_text.replace("<!-- image -->", "")
            
            if len(sheet_names) > 1:
                # Split content and insert sheet headers inline
                content_blocks = markdown_text.split("\n\n")
                blocks_per_sheet = len(content_blocks) // len(sheet_names)
                
                restructured_md = f"# {Path(file_path).name}\n\n"
                block_index = 0
                
                for sheet_idx, sheet_name in enumerate(sheet_names):
                    restructured_md += f"## {sheet_name}\n\n"
                    
                    if sheet_idx == len(sheet_names) - 1:
                        sheet_blocks = content_blocks[block_index:]
                    else:
                        sheet_blocks = content_blocks[block_index:block_index + blocks_per_sheet]
                    
                    restructured_md += "\n\n".join(block for block in sheet_blocks if block.strip())
                    restructured_md += "\n\n"
                    
                    block_index += blocks_per_sheet
                
                markdown_text = restructured_md
            elif len(sheet_names) == 1:
                # Single sheet - simpler header structure
                markdown_text = f"# {Path(file_path).name}\n\n## {sheet_names[0]}\n\n" + markdown_text
        except Exception as e:
            print(f"[WARNING] Could not inject sheet names: {e}")
        
        print(f"[SUCCESS] Extracted {len(markdown_text)} characters")
        return markdown_text
        
    except ImportError:
        print("[ERROR] Docling not installed")
        return None
    except Exception as e:
        print(f"[ERROR] {e}")
        return None



# ------------------------------
# MAIN: Run comparison tests on both file types
# Extracts using both methods and saves outputs for comparison
# ------------------------------

if __name__ == "__main__":
    print("\n" + "="*80)
    print("PowerPoint/Excel Extraction Comparison: Docling vs Specialized Libraries")
    print("="*80)
    
    # Test files stored in _local_uploads directory
    pptx_file = "_local_uploads/student_living_accommodation_sample.pptx" 
    excel_file = "_local_uploads/student_living_accommodation_sample.xlsx"
    
    # PowerPoint Extraction Comparison
    print("\n" + "-"*80)
    print("PowerPoint Extraction")
    print("-"*80)
    
    if Path(pptx_file).exists():
        docling_pptx = test_docling_pptx(pptx_file)
        if docling_pptx:
            Path("_local_uploads/extracted_pptx_docling.txt").write_text(docling_pptx, encoding='utf-8')
            print("✅ Saved: extracted_pptx_docling.txt")
        
        pptx_text = explore_pptx_extraction(pptx_file)
        if pptx_text:
            Path("_local_uploads/extracted_pptx_pythonpptx.txt").write_text(pptx_text, encoding='utf-8')
            print("✅ Saved: extracted_pptx_pythonpptx.txt")
    else:
        print(f"❌ File not found: {pptx_file}")
    
    # Excel Extraction Comparison
    print("\n" + "-"*80)
    print("Excel Extraction")
    print("-"*80)
    
    if Path(excel_file).exists():
        docling_excel = test_docling_excel(excel_file)
        if docling_excel:
            Path("_local_uploads/extracted_excel_docling.txt").write_text(docling_excel, encoding='utf-8')
            print("✅ Saved: extracted_excel_docling.txt")
        
        excel_text = explore_excel_extraction(excel_file)
        if excel_text:
            Path("_local_uploads/extracted_excel_openpyxl.txt").write_text(excel_text, encoding='utf-8')
            print("✅ Saved: extracted_excel_openpyxl.txt")
    else:
        print(f"❌ File not found: {excel_file}")
    
    print("\n" + "="*80)
    print("✅ Extraction complete. Compare extracted_*.txt files in _local_uploads/")
    print("="*80)

