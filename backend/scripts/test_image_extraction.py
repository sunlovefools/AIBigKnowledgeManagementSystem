"""
Test image extraction from PowerPoint and Excel files.
"""

import io
from pathlib import Path
import sys

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

def test_pptx_image_extraction():
    """Test extracting images from PowerPoint file using python-pptx."""
    print("=" * 80)
    print("TEST: PowerPoint Image Extraction")
    print("=" * 80)
    
    pptx_path = Path(__file__).resolve().parent.parent / "_local_uploads" / "student_living_accommodation_sample_WITH_TABLE.pptx"
    
    if not pptx_path.exists():
        print(f"❌ File not found: {pptx_path}")
        return False
    
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        with open(pptx_path, "rb") as f:
            prs = Presentation(io.BytesIO(f.read()))
        
        total_images = 0
        print(f"\n File: {pptx_path.name}")
        print(f" Total slides: {len(prs.slides)}\n")
        
        for slide_no, slide in enumerate(prs.slides, start=1):
            slide_images = 0
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_images += 1
                    total_images += 1
                    print(f"   Slide {slide_no}: Found picture (total: {total_images})")
            
            if slide_images == 0:
                print(f"   Slide {slide_no}: No images")
        
        print(f"\n✅ Total images found: {total_images}")
        
        if total_images == 0:
            print("   ⚠️  This PowerPoint file contains no images.")
        
        return True
        
    except ImportError:
        print("❌ python-pptx not installed. Install with: pip install python-pptx")
        return False
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_xlsx_image_extraction():
    """Test extracting images from Excel file using openpyxl."""
    print("\n" + "=" * 80)
    print("TEST: Excel Image Extraction")
    print("=" * 80)
    
    xlsx_path = Path(__file__).resolve().parent.parent / "_local_uploads" / "student_living_accommodation_sample.xlsx"
    
    if not xlsx_path.exists():
        print(f"❌ File not found: {xlsx_path}")
        return False
    
    try:
        import openpyxl
        
        with open(xlsx_path, "rb") as f:
            wb = openpyxl.load_workbook(io.BytesIO(f.read()))
        
        total_images = 0
        print(f"\n File: {xlsx_path.name}")
        print(f" Total sheets: {len(wb.sheetnames)}\n")
        
        for sheet_no, sheet_name in enumerate(wb.sheetnames, start=1):
            sheet = wb[sheet_name]
            sheet_images = 0
            
            if hasattr(sheet, '_images') and sheet._images:
                sheet_images = len(sheet._images)
                total_images += sheet_images
                print(f"   Sheet {sheet_no} ({sheet_name}): {sheet_images} image(s)")
            else:
                print(f"   Sheet {sheet_no} ({sheet_name}): No images")
        
        print(f"\n✅ Total images found: {total_images}")
        
        if total_images == 0:
            print("   ⚠️  This Excel file contains no images.")
        
        return True
        
    except ImportError:
        print("❌ openpyxl not installed. Install with: pip install openpyxl")
        return False
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_pptexcel_extractor_with_images():
    """Test the docling_pptexcel_extractor with files."""
    print("\n" + "=" * 80)
    print("TEST: Docling PPTX/XLSX Extractor (Full Pipeline)")
    print("=" * 80)
    
    from app.service.rag.ingestion.docling_pptexcel_extractor import parse_pptexcel_with_docling
    
    # Test PowerPoint
    pptx_path = Path(__file__).resolve().parent.parent / "_local_uploads" / "student_living_accommodation_sample_WITH_TABLE.pptx"
    
    if pptx_path.exists():
        print(f"\n Processing: {pptx_path.name}")
        
        with open(pptx_path, "rb") as f:
            file_bytes = f.read()
        
        try:
            result = parse_pptexcel_with_docling(
                file_bytes=file_bytes,
                file_name=pptx_path.name,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                file_id="test-pptx-001",
            )
            
            print(f"✅ Extraction successful!")
            print(f"   Structured blocks: {len(result.structured_blocks)}")
            print(f"   Images extracted: {len(result.images)}")
            print(f"   Warnings: {len(result.warnings)}")
            
            if result.images:
                print(f"\n    Image details:")
                for i, img in enumerate(result.images, start=1):
                    print(f"      {i}. UUID: {img.image_uuid}")
                    print(f"         Page: {img.page_no}, Index: {img.picture_index}")
                    print(f"         File: {img.file_name}")
                    print(f"         S3 Status: {img.s3_upload_status}")
            else:
                print(f"\n  No images found in this PowerPoint file.")
            
            # Check for image markers in blocks
            image_blocks = [b for b in result.structured_blocks if b.block_type == "picture"]
            if image_blocks:
                print(f"\n   Picture blocks created: {len(image_blocks)}")
            
        except Exception as e:
            print(f"❌ Extraction failed: {e}")
            import traceback
            traceback.print_exc()


if __name__ == "__main__":
    print("\n" + "╔" + "=" * 78 + "╗")
    print("║" + " " * 20 + "PPTX/XLSX IMAGE EXTRACTION TEST" + " " * 27 + "║")
    print("╙" + "=" * 78 + "╝")
    
    # Test raw extraction
    test_pptx_image_extraction()
    test_xlsx_image_extraction()
    
    # Test full pipeline
    test_pptexcel_extractor_with_images()
    
    print("\n" + "=" * 80)
    print("TESTS COMPLETE")
    print("=" * 80)
