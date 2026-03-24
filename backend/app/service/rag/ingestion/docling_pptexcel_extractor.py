"""
Docling-based extraction for PowerPoint and Excel files (PPTX/XLSX).

Purpose:
- Unified extraction for PPTX and XLSX files using Docling
- Returns structured blocks compatible with the Docling chunking pipeline
- Follows the same design pattern as docling_pdf_extractor.py
- Extracts images from PPTX/XLSX documents (similar to PDF image extraction)

Supported formats:
- PowerPoint (.pptx): application/vnd.openxmlformats-officedocument.presentationml.presentation
- Excel (.xlsx): application/vnd.openxmlformats-officedocument.spreadsheetml.sheet

Flow:
1. Convert PPTX/XLSX to Docling Document using DocumentConverter (text/tables)
2. Extract images separately using python-pptx / openpyxl (Docling limitation workaround)
3. Export document to markdown (structured text format)
4. Parse markdown into structured blocks for chunking
5. Inject image blocks with UUID markers (like PDF pipeline)
6. Return structured blocks compatible with docling_chunker.py
"""

from __future__ import annotations

import io
import re
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

from app.core.id_utils import generate_uuid_v6
from app.service.rag.ingestion.docling import table_image_vlm
from app.service.rag.ingestion.docling.models import (
    DoclingStructuredBlock,
    ExtractedImageArtifact,
)
from app.service.rag.ingestion.docling.storage import local_artifacts_store, s3_upload
from app.service.rag.ingestion.docling.utils import markdown_builder
from app.service.rag.ingestion.docling.utils.table_data_artifacts import (
    persist_table_data_toon_artifacts,
)
from app.service.rag.ingestion.markdown_canonicalizer import (
    canonicalize_docling_block_text,
)

# MIME types for PowerPoint/Excel documents
SUPPORTED_PPTEXCEL_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
}


class DoclingPptExcelParseResult:
    """
    Result from parsing a PowerPoint or Excel file with Docling.
    
    Similar to DoclingParseResult but simplified for PPTX/XLSX documents.
    Now includes image extraction like PDF pipeline.
    """
    
    def __init__(
        self,
        structured_blocks: list[DoclingStructuredBlock],
        markdown_content: str,
        warnings: list[str],
        file_id: str,
        file_name: str,
        images: list[ExtractedImageArtifact] | None = None,
    ):
        self.structured_blocks = structured_blocks
        self.markdown_content = markdown_content
        self.warnings = warnings
        self.file_id = file_id
        self.file_name = file_name
        self.images = images or []
        self.artifact_dir: Path | None = None
        self.artifact_run_id: str | None = None
        self.partial_failures: list[str] = []


def _extract_images_from_pptx(
    file_bytes: bytes,
    artifact_dir: Path | None,
    file_name: str,
    file_id: str,
    warnings: list[str],
) -> tuple[list[ExtractedImageArtifact], dict[int, list[str]], dict[str, int]]:
    """
    Extract images from PowerPoint file using python-pptx.

    Steps:
    1. Validate artifact destination and optional dependencies.
    2. Iterate slides and extract image bytes from picture shapes.
    3. Persist images, upload to S3, and collect warnings/statistics.
    """
    # 1) Initialize return containers for extracted artifacts and upload stats.
    images: list[ExtractedImageArtifact] = []
    slide_images: dict[int, list[str]] = {}  # slide_no -> [image_uuids]
    s3_stats = {"failed": 0, "uploaded": 0, "skipped": 0}

    # 2) Skip extraction when artifact output is disabled.
    if not artifact_dir:
        return images, slide_images, s3_stats

    # 3) Load optional PowerPoint dependency.
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        warnings.append(
            "python-pptx not installed - cannot extract PowerPoint images. "
            "Install with: pip install python-pptx"
        )
        return images, slide_images, s3_stats

    # 4) Iterate slides/shapes, persist images, and track S3 upload outcomes.
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        picture_counter = 0
        
        for slide_no, slide in enumerate(prs.slides, start=1):
            slide_image_uuids: list[str] = []
            
            for shape in slide.shapes:
                # Check if shape is a picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_counter += 1
                    
                    try:
                        # Extract image bytes
                        image_bytes = shape.image.blob
                        
                        # Generate UUID for this image
                        image_uuid = generate_uuid_v6()
                        
                        # Save image to artifact directory
                        image_name = f"{image_uuid}.png"
                        image_path = artifact_dir / "images" / image_name
                        image_path.parent.mkdir(parents=True, exist_ok=True)
                        image_path.write_bytes(image_bytes)
                        
                        # Create artifact
                        artifact = ExtractedImageArtifact(
                            kind="picture",
                            image_uuid=image_uuid,
                            file_name=image_name,
                            file_path=str(image_path),
                            page_no=slide_no,
                            picture_index=picture_counter,
                        )
                        
                        # Upload to S3 (same as PDF pipeline)
                        artifact = s3_upload.upload_image_artifact_to_s3(
                            artifact,
                            source_file_name=file_name,
                            file_id=file_id,
                        )
                        
                        if artifact.s3_upload_status == "failed":
                            s3_stats["failed"] += 1
                            warnings.append(
                                f"Failed to upload PowerPoint image_uuid={artifact.image_uuid} to S3: {artifact.s3_error}"
                            )
                        elif artifact.s3_upload_status == "uploaded":
                            s3_stats["uploaded"] += 1
                        elif artifact.s3_upload_status == "skipped":
                            s3_stats["skipped"] += 1
                        
                        images.append(artifact)
                        slide_image_uuids.append(image_uuid)
                        
                    except Exception as e:
                        warnings.append(
                            f"Failed to extract image #{picture_counter} from slide {slide_no}: {e}"
                        )
            
            if slide_image_uuids:
                slide_images[slide_no] = slide_image_uuids
        
        print(f"[docling-pptexcel] Extracted {len(images)} images from PowerPoint")
        print(f"[docling-pptexcel] S3 upload: {s3_stats['uploaded']} uploaded, {s3_stats['skipped']} skipped, {s3_stats['failed']} failed")
        
    except Exception as e:
        warnings.append(f"Failed to extract PowerPoint images: {e}")

    return images, slide_images, s3_stats


def _extract_images_from_xlsx(
    file_bytes: bytes,
    artifact_dir: Path | None,
    file_name: str,
    file_id: str,
    warnings: list[str],
) -> tuple[list[ExtractedImageArtifact], dict[int, list[str]], dict[str, int]]:
    """
    Extract images from Excel file using openpyxl.

    Steps:
    1. Validate artifact destination and optional dependencies.
    2. Iterate worksheets and extract embedded image bytes.
    3. Persist images, upload to S3, and collect warnings/statistics.
    """
    # 1) Initialize return containers for extracted artifacts and upload stats.
    images: list[ExtractedImageArtifact] = []
    sheet_images: dict[int, list[str]] = {}  # sheet_no -> [image_uuids]
    s3_stats = {"failed": 0, "uploaded": 0, "skipped": 0}

    # 2) Skip extraction when artifact output is disabled.
    if not artifact_dir:
        return images, sheet_images, s3_stats

    # 3) Load optional Excel dependency.
    try:
        import openpyxl
    except ImportError:
        warnings.append(
            "openpyxl not installed - cannot extract Excel images. "
            "Install with: pip install openpyxl"
        )
        return images, sheet_images, s3_stats

    # 4) Iterate workbook sheets, persist images, and track S3 upload outcomes.
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
        picture_counter = 0
        
        for sheet_no, sheet_name in enumerate(wb.sheetnames, start=1):
            sheet = wb[sheet_name]
            sheet_image_uuids: list[str] = []
            
            # Check if sheet has images
            if hasattr(sheet, '_images') and sheet._images:
                for img in sheet._images:
                    picture_counter += 1
                    
                    try:
                        # Extract image bytes
                        image_bytes = img._data()
                        
                        # Generate UUID for this image
                        image_uuid = generate_uuid_v6()
                        
                        # Save image to artifact directory
                        image_name = f"{image_uuid}.png"
                        image_path = artifact_dir / "images" / image_name
                        image_path.parent.mkdir(parents=True, exist_ok=True)
                        image_path.write_bytes(image_bytes)
                        
                        # Create artifact
                        artifact = ExtractedImageArtifact(
                            kind="picture",
                            image_uuid=image_uuid,
                            file_name=image_name,
                            file_path=str(image_path),
                            page_no=sheet_no,
                            picture_index=picture_counter,
                        )
                        
                        # Upload to S3 (same as PDF pipeline)
                        artifact = s3_upload.upload_image_artifact_to_s3(
                            artifact,
                            source_file_name=file_name,
                            file_id=file_id,
                        )
                        
                        if artifact.s3_upload_status == "failed":
                            s3_stats["failed"] += 1
                            warnings.append(
                                f"Failed to upload Excel image_uuid={artifact.image_uuid} to S3: {artifact.s3_error}"
                            )
                        elif artifact.s3_upload_status == "uploaded":
                            s3_stats["uploaded"] += 1
                        elif artifact.s3_upload_status == "skipped":
                            s3_stats["skipped"] += 1
                        
                        images.append(artifact)
                        sheet_image_uuids.append(image_uuid)
                        
                    except Exception as e:
                        warnings.append(
                            f"Failed to extract image #{picture_counter} from sheet '{sheet_name}': {e}"
                        )
            
            if sheet_image_uuids:
                sheet_images[sheet_no] = sheet_image_uuids
        
        print(f"[docling-pptexcel] Extracted {len(images)} images from Excel")
        print(f"[docling-pptexcel] S3 upload: {s3_stats['uploaded']} uploaded, {s3_stats['skipped']} skipped, {s3_stats['failed']} failed")
        
    except Exception as e:
        warnings.append(f"Failed to extract Excel images: {e}")

    return images, sheet_images, s3_stats


def _might_contain_table(image_artifact: ExtractedImageArtifact) -> bool:
    """
    Heuristic to determine if an image might contain a table.

    Steps:
    1. Read image dimensions when available.
    2. Apply a size threshold heuristic.
    3. Fail open (True) if analysis cannot run.
    """
    # 1) Attempt to inspect image dimensions using PIL.
    try:
        from PIL import Image

        image_path = Path(image_artifact.file_path)
        if not image_path.exists():
            return False

        # 2) Use area threshold as a conservative table-likelihood signal.
        with Image.open(image_path) as img:
            width, height = img.size
            return (width * height) > 50000

    # 3) Fail open so potential tables are not silently missed.
    except Exception:
        return True


def is_pptexcel_document(content_type: str) -> bool:
    """
    Check if a content type is a supported PowerPoint/Excel document format.

    Steps:
    1. Receive MIME type.
    2. Compare against supported PPTX/XLSX MIME constants.
    3. Return boolean match result.
    """
    return content_type in SUPPORTED_PPTEXCEL_MIME_TYPES


def _build_pipeline_markdown_parts(
    markdown_content: str,
    file_name: str,
    page_images: dict[int, list[str]] | None = None,
    *,
    artifact_dir: Path | None,
    images_by_uuid: dict[str, ExtractedImageArtifact],
    table_image_candidate_uuids: set[str],
    table_image_vlm_runtime: Any | None,
) -> tuple[list[str], list[dict[str, Any]], list[table_image_vlm.TableImageVlmJob]]:
    """
    Parse Office markdown into pipeline-style markdown parts and metadata.

    Steps:
    1. Segment markdown text into typed blocks (header/text/list/table).
    2. Inject image/table-image marker blocks by logical page/sheet.
    3. Emit markdown parts and metadata aligned with markdown_builder contract.
    """

    # 1) Initialize mutable accumulators used by parser and VLM queueing.
    lines = markdown_content.split("\n")
    markdown_parts: list[str] = []
    structured_block_metadata: list[dict[str, Any]] = []
    table_image_vlm_jobs: list[table_image_vlm.TableImageVlmJob] = []

    current_block_lines: list[str] = []
    current_block_type = "text"
    current_page_no = 1
    page_images = page_images or {}
    injected_pages: set[int] = set()
    table_index = 0

    slide_number_pattern = re.compile(r"^##\s+Slide\s+(\d+)", re.IGNORECASE)
    sheet_pattern = re.compile(r"^##\s+Sheet:\s+(.+)", re.IGNORECASE)

    def _flush_current_block() -> None:
        """Finalize the current text/list/table/header block into markdown parts."""
        nonlocal current_block_lines

        if not current_block_lines:
            return

        content = "\n".join(current_block_lines).strip()
        current_block_lines = []
        if not content:
            return

        markdown_builder.append_markdown_block(
            markdown_parts=markdown_parts,
            structured_block_metadata=structured_block_metadata,
            text=content,
            block_type=current_block_type,
            page_no=current_page_no,
            is_table_image=False,
            table_image_uuid=None,
        )

    def _append_image_block(*, page_no: int, image_uuid: str) -> None:
        """Append picture block or table-image block and queue VLM job when eligible."""
        nonlocal table_index

        image_artifact = images_by_uuid.get(image_uuid)
        should_queue_vlm = (
            image_uuid in table_image_candidate_uuids
            and image_artifact is not None
            and artifact_dir is not None
            and table_image_vlm_runtime is not None
        )

        if should_queue_vlm:
            table_index += 1
            summary_placeholder = table_image_vlm.table_image_vlm_summary_placeholder(
                image_uuid
            )
            table_markdown_text = "\n".join(
                [
                    "> **Table (image)**: Table exists in image form.",
                    f"> {markdown_builder.table_image_uuid_marker(image_uuid)}",
                    f"> ![{image_artifact.file_name}]({local_artifacts_store.image_markdown_rel_path_from_uuid(image_uuid)})",
                    f"> {summary_placeholder}",
                ]
            )
            markdown_builder.append_markdown_block(
                markdown_parts=markdown_parts,
                structured_block_metadata=structured_block_metadata,
                text=table_markdown_text,
                block_type="table",
                page_no=page_no,
                is_table_image=True,
                table_image_uuid=image_uuid,
            )
            table_image_vlm_jobs.append(
                table_image_vlm.TableImageVlmJob(
                    image_artifact=image_artifact,
                    table_index=table_index,
                    page_no=page_no,
                    block_index=len(markdown_parts) - 1,
                    summary_placeholder=summary_placeholder,
                    output_dir=table_image_vlm.table_image_vlm_output_dir(
                        artifact_dir,
                        table_index=table_index,
                        image_uuid=image_uuid,
                    ),
                    json_rel_path=table_image_vlm.table_image_vlm_json_rel_path(
                        table_index=table_index,
                        image_uuid=image_uuid,
                    ),
                )
            )
            return

        markdown_builder.append_markdown_block(
            markdown_parts=markdown_parts,
            structured_block_metadata=structured_block_metadata,
            text=markdown_builder.picture_uuid_marker(image_uuid),
            block_type="picture",
            page_no=page_no,
            is_table_image=False,
            table_image_uuid=None,
        )

    def _inject_image_blocks_for_page(page_no: int) -> None:
        """Inject image-derived blocks once per logical page/sheet."""
        if page_no in injected_pages:
            return
        injected_pages.add(page_no)
        for image_uuid in page_images.get(page_no, []):
            _append_image_block(page_no=page_no, image_uuid=image_uuid)

    # 2) Parse markdown into typed blocks while preserving simple heuristics.
    for line in lines:
        stripped_line = line.strip()

        if not stripped_line and not current_block_lines:
            continue

        if stripped_line.startswith("#"):
            _flush_current_block()

            slide_match = slide_number_pattern.match(stripped_line)
            if slide_match:
                current_page_no = int(slide_match.group(1))

            sheet_match = sheet_pattern.match(stripped_line)
            if sheet_match:
                current_page_no += 1

            current_block_type = "header"
            current_block_lines.append(stripped_line.lstrip("#").strip())
            _flush_current_block()
            _inject_image_blocks_for_page(current_page_no)
            continue

        if "|" in stripped_line and stripped_line.count("|") >= 2:
            if current_block_type != "table":
                _flush_current_block()
                current_block_type = "table"
            if re.match(r"^\|[\s\-:]+\|", stripped_line):
                continue
            current_block_lines.append(stripped_line)
            continue

        if stripped_line.startswith(("-", "*")) or re.match(r"^\d+\.\s", stripped_line):
            if current_block_type != "list":
                _flush_current_block()
                current_block_type = "list"
            current_block_lines.append(stripped_line)
            continue

        if current_block_type not in {"text", "list", "table"}:
            _flush_current_block()
            current_block_type = "text"

        if not stripped_line and current_block_type in {"table", "list"}:
            _flush_current_block()
            current_block_type = "text"

        if stripped_line or current_block_lines:
            current_block_lines.append(line)

    # 3) Flush trailing block and inject any remaining page/sheet images.
    _flush_current_block()
    for page_no in sorted(set(page_images.keys())):
        _inject_image_blocks_for_page(page_no)

    _ = file_name
    print(
        "[docling-pptexcel] Parsed %s markdown parts and queued %s VLM job(s)"
        % (len(markdown_parts), len(table_image_vlm_jobs))
    )
    return markdown_parts, structured_block_metadata, table_image_vlm_jobs


def parse_pptexcel_with_docling(
    file_bytes: bytes,
    file_name: str,
    content_type: str,
    file_id: str | None = None,
) -> DoclingPptExcelParseResult:
    """
    Parse PowerPoint/Excel file with Docling and extract text/tables/images.

    Steps:
    1. Convert Office document bytes to Docling markdown output.
    2. Extract Office images and build pipeline-style markdown parts/metadata.
    3. Finalize VLM jobs, canonicalize blocks, and emit parse artifacts/result.
    """
    # 1) Validate request payload and MIME type before any heavy processing.
    if not file_bytes:
        raise ValueError("empty file payload")

    if not is_pptexcel_document(content_type):
        raise ValueError(
            f"Unsupported content type for Docling PPTX/XLSX extraction: {content_type}. "
            f"Supported types: {', '.join(SUPPORTED_PPTEXCEL_MIME_TYPES)}"
        )

    # 2) Resolve file-type labels and shared parse state.
    file_type = "PowerPoint" if "presentation" in content_type else "Excel"
    print(
        f"[docling-pptexcel] Processing {file_type} file: {file_name} ({len(file_bytes)} bytes)"
    )

    resolved_file_id = (file_id or generate_uuid_v6()).strip()
    warnings: list[str] = []

    # 3) Load Docling converter lazily to avoid startup dependency cost.
    try:
        from docling.document_converter import DocumentConverter
    except ImportError as e:
        raise ImportError(
            f"Docling library not installed. Install with: pip install docling. Error: {e}"
        )

    # 4) Execute full Office parse pipeline and materialize structured outputs.
    try:
        print(f"[docling-pptexcel] Converting {file_type} with Docling...")
        import tempfile
        import os

        extension = Path(file_name).suffix
        with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name

        try:
            converter = DocumentConverter()
            result = converter.convert(temp_file_path)
            print("[docling-pptexcel] Conversion successful")
        finally:
            try:
                os.unlink(temp_file_path)
            except Exception as e:
                print(
                    f"[docling-pptexcel] Warning: Could not delete temp file {temp_file_path}: {e}"
                )

        print("[docling-pptexcel] Exporting to markdown...")
        markdown_content = result.document.export_to_markdown()
        print(
            f"[docling-pptexcel] Exported {len(markdown_content)} characters of markdown"
        )

        # 4.1) Add stable file/sheet headers so downstream chunking has structure anchors.
        if file_type == "Excel":
            try:
                print("[docling-pptexcel] Adding Excel sheet names as headers...")
                sheet_names: list[str] = []
                groups = list(getattr(result.document, "groups", []) or [])
                for i, grp in enumerate(groups):
                    raw_name = str(getattr(grp, "name", "") or "").strip()
                    if not raw_name:
                        continue
                    # Expected format from Docling groups: "sheet: <name>".
                    if ":" in raw_name:
                        prefix, suffix = raw_name.split(":", 1)
                        if prefix.strip().lower() == "sheet":
                            raw_name = suffix.strip()
                    cleaned_name = raw_name.strip() or f"Sheet {i + 1}"
                    sheet_names.append(cleaned_name)

                if not sheet_names:
                    sheet_names = ["Sheet 1"]

                num_sheets = len(sheet_names)
                markdown_content = markdown_content.replace("<!-- image -->", "")

                if num_sheets > 1:
                    content_blocks = markdown_content.split("\n\n")
                    blocks_per_sheet = len(content_blocks) // num_sheets if num_sheets else 0
                    restructured_content = f"# {file_name}\n\n"
                    block_index = 0
                    for sheet_idx, sheet_name in enumerate(sheet_names):
                        restructured_content += f"## {sheet_name}\n\n"
                        if sheet_idx == len(sheet_names) - 1:
                            sheet_blocks = content_blocks[block_index:]
                        else:
                            sheet_blocks = content_blocks[
                                block_index : block_index + blocks_per_sheet
                            ]
                        restructured_content += "\n\n".join(
                            block for block in sheet_blocks if block.strip()
                        )
                        restructured_content += "\n\n"
                        block_index += blocks_per_sheet
                    markdown_content = restructured_content
                    print(
                        f"[docling-pptexcel] Restructured markdown with headers for {num_sheets} Excel sheets: {sheet_names}"
                    )
                else:
                    markdown_content = (
                        f"# {file_name}\n\n## {sheet_names[0]}\n\n" + markdown_content
                    )
                    print(
                        f"[docling-pptexcel] Added headers for single Excel sheet: {sheet_names[0]}"
                    )
            except Exception as e:
                warnings.append(f"Could not add Excel sheet names to markdown: {e}")
                print(f"[docling-pptexcel] Warning: Could not add sheet names: {e}")
        else:
            markdown_content = markdown_content.replace("<!-- image -->", "")
            markdown_content = f"# {file_name}\n\n" + markdown_content
            print("[docling-pptexcel] Added filename header for PowerPoint")

        # 4.2) Prepare artifacts and extract Office images.
        print(f"[docling-pptexcel] Extracting images from {file_type}...")
        run_id, artifact_dir, markdown_path = local_artifacts_store.prepare_docling_artifact_dir(
            file_name=file_name,
            artifact_root=None,
        )

        if file_type == "PowerPoint":
            images, page_images, _s3_stats = _extract_images_from_pptx(
                file_bytes=file_bytes,
                artifact_dir=artifact_dir,
                file_name=file_name,
                file_id=resolved_file_id,
                warnings=warnings,
            )
        else:
            images, page_images, _s3_stats = _extract_images_from_xlsx(
                file_bytes=file_bytes,
                artifact_dir=artifact_dir,
                file_name=file_name,
                file_id=resolved_file_id,
                warnings=warnings,
            )

        print(f"[docling-pptexcel] Extracted {len(images)} images from {file_type}")

        # 4.3) Prepare VLM runtime and candidate table-image UUID set.
        table_image_vlm_runtime = (
            table_image_vlm.build_table_image_vlm_runtime(
                artifact_dir=artifact_dir,
                warnings=warnings,
            )
            if artifact_dir is not None
            else None
        )
        table_image_candidate_uuids = {
            image_artifact.image_uuid
            for image_artifact in images
            if _might_contain_table(image_artifact)
        }
        if table_image_vlm_runtime is None:
            print("[docling-pptexcel] Table image VLM disabled (check OPENROUTER_API_KEY)")
        elif not table_image_candidate_uuids:
            print("[docling-pptexcel] No images detected as potential tables")
        else:
            print(
                "[docling-pptexcel] Table image VLM enabled - %s candidate image(s)"
                % len(table_image_candidate_uuids)
            )

        # 4.4) Build pipeline-style markdown parts/metadata (chunker-compatible).
        (
            markdown_parts,
            structured_block_metadata,
            table_image_vlm_jobs,
        ) = _build_pipeline_markdown_parts(
            markdown_content=markdown_content,
            file_name=file_name,
            page_images=page_images,
            artifact_dir=artifact_dir,
            images_by_uuid={image.image_uuid: image for image in images},
            table_image_candidate_uuids=table_image_candidate_uuids,
            table_image_vlm_runtime=table_image_vlm_runtime,
        )

        # 4.5) Submit/finalize shared table-image VLM jobs and persist TOON artifacts.
        table_image_vlm_executor: ThreadPoolExecutor | None = None
        if (
            table_image_vlm_runtime is not None
            and artifact_dir is not None
            and table_image_vlm_jobs
        ):
            table_image_vlm_executor = ThreadPoolExecutor(
                max_workers=table_image_vlm_runtime.max_workers,
                thread_name_prefix="table-vlm-pptexcel",
            )
        try:
            if (
                table_image_vlm_runtime is not None
                and table_image_vlm_executor is not None
                and artifact_dir is not None
                and table_image_vlm_jobs
            ):
                table_image_vlm.submit_ready_table_image_vlm_jobs(
                    runtime=table_image_vlm_runtime,
                    executor=table_image_vlm_executor,
                    jobs=table_image_vlm_jobs,
                    markdown_parts=markdown_parts,
                    warnings=warnings,
                    force=True,
                )
                table_image_vlm.finalize_table_image_vlm_jobs(
                    artifact_dir=artifact_dir,
                    jobs=table_image_vlm_jobs,
                    markdown_parts=markdown_parts,
                    warnings=warnings,
                )
                persist_table_data_toon_artifacts(
                    artifact_dir=artifact_dir,
                    table_image_vlm_jobs=table_image_vlm_jobs,
                    resolved_file_id=resolved_file_id,
                    file_name=file_name,
                    warnings=warnings,
                )
        finally:
            if table_image_vlm_executor is not None:
                table_image_vlm_executor.shutdown(wait=True)

        # 4.6) Canonicalize markdown by block type and materialize structured blocks.
        canonicalized_markdown_parts = list(markdown_parts)
        for metadata in structured_block_metadata:
            block_index = int(metadata.get("block_index", -1))
            if block_index < 0 or block_index >= len(canonicalized_markdown_parts):
                continue
            block_type = str(metadata.get("block_type") or "text")
            canonicalized_markdown_parts[block_index] = canonicalize_docling_block_text(
                block_type=block_type,
                text=canonicalized_markdown_parts[block_index],
            )

        markdown_parts = canonicalized_markdown_parts
        markdown_text = "\n\n".join(markdown_parts)
        structured_blocks = markdown_builder.build_structured_blocks(
            structured_block_metadata=structured_block_metadata,
            markdown_parts=markdown_parts,
        )
        if not structured_blocks:
            warnings.append(
                f"Docling produced markdown but no structured blocks were extracted for {file_name}"
            )

        # 4.7) Persist markdown + manifest artifacts when enabled.
        if artifact_dir is not None and markdown_path is not None:
            markdown_path.write_text(markdown_text, encoding="utf-8")

        # 4.8) Build and return stable Office parse result contract.
        parse_result = DoclingPptExcelParseResult(
            structured_blocks=structured_blocks,
            markdown_content=markdown_text,
            warnings=warnings,
            file_id=resolved_file_id,
            file_name=file_name,
            images=images,
        )
        parse_result.artifact_dir = artifact_dir
        parse_result.artifact_run_id = run_id or ""

        if artifact_dir is not None and markdown_path is not None:
            local_artifacts_store.write_manifest(
                artifact_dir,
                {
                    "source_file_name": file_name,
                    "artifact_run_id": run_id or "",
                    "artifact_dir": str(artifact_dir),
                    "markdown_path": str(markdown_path),
                    "markdown_text": markdown_text,
                    "images": [image.model_dump() for image in images],
                    "warnings": warnings,
                    "partial_failures": parse_result.partial_failures,
                    "stats": {
                        "converted_chunks": 1,
                        "partial_failure_chunks": 0,
                        "pictures_extracted": sum(
                            1 for image in images if image.kind == "picture"
                        ),
                        "table_fallback_images_extracted": len(table_image_candidate_uuids),
                    },
                    "structured_blocks": [
                        block.model_dump() for block in structured_blocks
                    ],
                },
            )

        return parse_result

    except Exception as exc:
        error_msg = f"Docling {file_type} processing failed for {file_name}: {exc}"
        print(f"[docling-pptexcel] ERROR: {error_msg}")
        raise Exception(error_msg) from exc

def get_pptexcel_ingestion_strategy() -> str:
    """
    Determine the PowerPoint/Excel document ingestion strategy.

    Steps:
    1. Keep a dedicated strategy hook for Office ingestion.
    2. Mirror the PDF strategy-pattern call shape.
    3. Return fixed `"docling"` for current Office flow.
    """
    return "docling"


__all__ = [
    "SUPPORTED_PPTEXCEL_MIME_TYPES",
    "DoclingPptExcelParseResult",
    "is_pptexcel_document",
    "parse_pptexcel_with_docling",
    "get_pptexcel_ingestion_strategy",
]
